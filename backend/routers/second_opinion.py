"""IA Contre-analyse (Second Opinion) : POST /api/second-opinion/*.

On soumet un DOCUMENT EXTERNE (argumentaire de broker en PDF, note de conseil en
PPTX) recu par l'investisseur, plus ses consignes ; le modele produit une
contre-analyse en confrontant le document aux donnees Barzel de la ville.

Deux endpoints :
- POST /extract (multipart) : extrait le texte des fichiers (PDF via pypdf, PPTX
  via python-pptx). Le front stocke ce texte dans la conversation (localStorage).
- POST /analyze (JSON) : recoit le texte du document + l'historique metier de la
  conversation (continuation consciente, backend sans etat) et renvoie l'analyse.

SECURITE : le contexte Barzel et les garde-fous sont EXACTEMENT ceux de l'analyste
(memes fonctions importees, aucune duplication), plus le texte du document fourni
par l'utilisateur. Jamais de simulation, jamais params.json, maille par pays, zero
tiret cadratin (filet strip_em_dashes). Cle Anthropic lue de backend/.env.
"""

from __future__ import annotations

import io
import logging

from fastapi import APIRouter, File, HTTPException, UploadFile
from pydantic import BaseModel

from ..services.cities import label_for
from .analyst import (
    _MODEL,
    _LANG_NAME,
    _VERDICT_VOCAB,
    _api_key,
    _build_context,
    _norm_lang,
    _require_city,
    _require_class,
    mesh_for,
    strip_em_dashes,
)

log = logging.getLogger("routers.second_opinion")

router = APIRouter(prefix="/api/second-opinion", tags=["second-opinion"])

_MAX_TOKENS = 8000  # analyse complete ; marge pour ne jamais tronquer la recommandation finale
_MAX_FILES = 3
_MAX_BYTES = 20 * 1024 * 1024  # 20 Mo par fichier
_MAX_DOC_CHARS = 60_000  # borne du texte extrait (envoi LLM + stockage front)
_CLS_ALL = "all"  # « toutes classes confondues » : dossier mixte, contexte des 5 classes
_CLS_ORDER = ("residential", "office", "hotel", "logistics", "retail")


# --------------------------------------------------------------------------- #
# Extraction : PDF (pypdf) + PPTX (python-pptx)                                #
# --------------------------------------------------------------------------- #

def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:  # noqa: BLE001 ; une page illisible ne casse pas l'extraction
            continue
    return "\n".join(parts).strip()


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        texts.append(" | ".join(cells))
        if texts:
            parts.append(f"[Diapositive {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts).strip()


def _extract_one(filename: str, data: bytes) -> str:
    name = (filename or "").lower()
    if name.endswith(".pdf"):
        return _extract_pdf(data)
    if name.endswith(".pptx"):
        return _extract_pptx(data)
    raise HTTPException(status_code=400, detail=f"format non supporte : {filename!r} (PDF ou PPTX)")


@router.post("/extract")
async def extract(files: list[UploadFile] = File(...)) -> dict:
    if not files:
        raise HTTPException(status_code=400, detail="aucun fichier")
    if len(files) > _MAX_FILES:
        raise HTTPException(status_code=400, detail=f"maximum {_MAX_FILES} fichiers")
    blocks: list[str] = []
    names: list[str] = []
    for f in files:
        data = await f.read()
        if len(data) > _MAX_BYTES:
            raise HTTPException(status_code=400, detail=f"fichier trop volumineux : {f.filename!r}")
        text = _extract_one(f.filename or "", data)
        names.append(f.filename or "document")
        blocks.append(f"### DOCUMENT : {f.filename}\n{text}" if text
                      else f"### DOCUMENT : {f.filename}\n(aucun texte extractible)")
    doc_text = "\n\n".join(blocks)[:_MAX_DOC_CHARS]
    return {"doc_text": doc_text, "doc_names": names, "chars": len(doc_text)}


# --------------------------------------------------------------------------- #
# Analyse : document + consignes croises au contexte Barzel                    #
# --------------------------------------------------------------------------- #

class Turn(BaseModel):
    role: str
    text: str


class AnalyzePayload(BaseModel):
    doc_text: str = ""
    messages: list[Turn] = []  # historique metier ; messages[0] = consignes de l'investisseur
    asset_class: str = "residential"
    city: str | None = None
    lang: str = "fr"


_SYSTEM_TEMPLATE = """Tu es l'analyste de Barzel Analytics, plateforme d'intelligence immobilière couvrant {VILLE} pour un investisseur institutionnel.

On te soumet un DOCUMENT EXTERNE (argumentaire d'un broker, note d'un conseil, dossier de vente) reçu par l'investisseur, accompagné de ses consignes (questions, intentions, recommandations). Ta mission : produire une CONTRE-ANALYSE pour un comité d'investissement, en confrontant les affirmations du document aux données Barzel de {VILLE}.

STRUCTURE DE LA NOTE, DANS CET ORDRE :
1. Une phrase de recommandation en tête, qui porte le verdict pertinent.
2. La passe fiscale ci-dessous, menée comme tout premier volet d'analyse, avant tout le reste.
3. Les autres constats, ordonnés par impact chiffré décroissant.
4. La section finale des points non tranchables.

PASSE FISCALE OBLIGATOIRE, EN PREMIER, NON NÉGOCIABLE : si le document présente un mix de typologies avec des surfaces et un prix au mètre carré, mène-la avant toute autre analyse.
- Portugal : calcule le prix de vente par logement pour chaque typologie ; compare-le au seuil de 660 982 euros ; indique combien d'unités passent au-dessus et combien restent en dessous ; calcule la surface pivot (660 982 divisé par le prix au mètre carré). Si le document affirme un régime de TVA, vérifie cette affirmation et chiffre l'écart de marge si elle est fausse.
- Belgique : pas de seuil de prix ; le régime bascule selon la destination, 6 % en location longue durée en résidence principale sous plafond de 200 mètres carrés habitables, 21 % en vente libre. Détermine ce que retient le document, vérifie l'éligibilité au plafond de surface, et chiffre l'écart de marge entre les deux régimes, en euros et en euros par mètre carré vendable.
Si le document ne permet pas de mener cette passe, dis-le explicitement dans la section des points non tranchables.

DISCIPLINE D'UNITÉ : toute grandeur au mètre carré précise sa base, surface brute de construction ou surface privative vendable ; vérifie que les unités restent cohérentes du début à la fin de la note.
- Base des données Barzel : les coûts sont rapportés au mètre carré de surface vendable, et le coût de revient total est un coût tout compris (il intègre déjà les frais annexes et le portage financier). Avant d'utiliser une valeur, pose sa base ; ne t'en sers jamais sans l'avoir identifiée.
- Conversion obligatoire : un coût que le document exprime sur une autre base (surface brute de construction, coût dur seul) doit être ramené à la base Barzel avant toute comparaison, et la conversion doit être montrée. Un coût dur seul doit d'abord être majoré des postes annexes équivalents (provision pour aléas, honoraires, frais financiers, commercialisation) avant d'être comparé à un coût tout compris.
- Ligne de gamme : retiens l'hypothèse de coût qui correspond au niveau de prestation décrit par le document, jamais la valeur par défaut. Un programme classé A, NZEB ou haut de gamme relève de la fourchette haute, pas de la ligne standard ; justifie ce choix en une phrase.
- Cohérence entre sources : quand plusieurs grandeurs Barzel couvrent le même poste, cite-les toutes et explique laquelle s'applique, plutôt que d'en retenir une sans le dire.

CONFRONTATION AU DOCUMENT, À APPLIQUER SYSTÉMATIQUEMENT :
- Critique la représentativité des comparables avancés (localisation, période, standing, taille de l'échantillon).
- Repère les postes absents du bilan promoteur, par exemple provision pour aléas, frais financiers, commercialisation, honoraires, infrastructures et VRD.
- Signale les régimes fiscaux d'acquisition que le document passe sous silence.
- Pour chaque chiffre avancé, mesure l'écart au chiffre Barzel : angle mort, risque non dit, hypothèse optimiste.

RÉCONCILIATION DES SIGNAUX : quand deux indicateurs Barzel pointent en sens opposés (par exemple une valeur foncière de marché supérieure à la valeur foncière résiduelle calculée), ne les empile pas ; pose la contradiction explicitement et propose la lecture qui la résout.

GRANDEURS NON COMPARABLES : n'oppose jamais un délai de cession d'un actif existant à un délai d'écoulement d'un programme neuf, ce sont deux grandeurs distinctes ; toute comparaison de délais porte sur des grandeurs de même nature. Le délai d'écoulement d'un programme neuf relève du rythme d'absorption du marché (le rythme de vente pour ce niveau de prix et ce type de produit), jamais du délai de cession d'un actif existant.

RÈGLES ABSOLUES :
- Ta référence chiffrée est la section DONNÉES BARZEL. Tu n'inventes JAMAIS un chiffre côté Barzel : chaque nombre Barzel cité figure tel quel dans les données. Les scores se citent en entiers (« 87/100 »).
- Tu peux reprendre les chiffres AVANCÉS PAR LE DOCUMENT, mais tu les attribues clairement à lui (« le broker avance ... », « le dossier annonce ... ») et tu les confrontes aux données Barzel.
- Tu nommes les {MESH_PL} concernées. Tu ne mentionnes JAMAIS de niveau de confiance, de source, de méthodologie interne, ni l'idée qu'une donnée Barzel serait simulée ou estimée.
- Si le document sort du périmètre de {VILLE} ou de l'immobilier couvert par la plateforme, dis-le avec élégance et traite ce qui est couvrable.
- REGISTRE MÉTIER : vocabulaire immobilier uniquement, aucun terme interne de la plateforme dans ta réponse ; jamais « mode landbank » (dis réserve foncière ou potentiel de constructibilité), jamais « mode arbitrage » (dis fenêtre de cession ou de revente), jamais « uplift » (dis valorisation foncière du changement d'usage) ; n'introduis pas les verdicts comme des artefacts de la plateforme (« le verdict En attente »), énonce-les comme ta conclusion.
- FRANÇAIS CORRECT quand tu réponds en français : aucun calque du portugais (« provision pour aléas », jamais « contingence d'oeuvre ») ; orthographie exactement les noms propres portugais.
- Ponctuation : JAMAIS de tiret cadratin (le tiret long, U+2014), ni seul ni encadré d'espaces ; articule avec deux-points, virgule, parenthèses ou une nouvelle phrase.
- Ton sobre et professionnel, registre de comité d'investissement. Analyse structurée en paragraphes (PAS de markdown : pas de titres, pas de gras, pas de puces), développée mais sans remplissage : concentre-toi sur les écarts et les risques qui comptent pour la décision.
- Réponds en priorité aux consignes explicites de l'investisseur.
- Rédige TOUTE ta réponse en {LANG_NAME}. La recommandation d'ouverture emploie EXACTEMENT le libellé de verdict pertinent, dans la langue de réponse : {VERDICT_VOCAB}.

SECTION FINALE OBLIGATOIRE : termine par une section listant ce qui ne peut pas être tranché sans données complémentaires du client, et, pour chaque point, la donnée précise qui serait nécessaire."""


def _system_for(city: str, lang: str) -> str:
    return (_SYSTEM_TEMPLATE
            .replace("{VILLE}", label_for(city))
            .replace("{MESH_PL}", mesh_for(city)[1])
            .replace("{LANG_NAME}", _LANG_NAME[lang])
            .replace("{VERDICT_VOCAB}", _VERDICT_VOCAB[lang]))


def _context_for(asset_class: str, city: str) -> str:
    """Contexte Barzel : une classe, ou les cinq concatenees si « all » (dossier
    mixte residentiel + commerce, etc.). Chaque bloc porte deja son en-tete de classe."""
    if asset_class == _CLS_ALL:
        return "\n\n".join(_build_context(c, city) for c in _CLS_ORDER)
    return _build_context(asset_class, city)


@router.post("/analyze")
def analyze(payload: AnalyzePayload) -> dict:
    raw_cls = payload.asset_class or "residential"
    asset_class = _CLS_ALL if raw_cls == _CLS_ALL else _require_class(raw_cls)
    city = _require_city(payload.city)
    lang = _norm_lang(payload.lang)
    doc_text = (payload.doc_text or "").strip()[:_MAX_DOC_CHARS]
    turns = payload.messages
    if not turns or turns[0].role != "user":
        raise HTTPException(status_code=400, detail="conversation vide ou mal formee")
    if not doc_text:
        raise HTTPException(status_code=400, detail="document manquant")

    key = _api_key()
    if not key:
        raise HTTPException(status_code=503, detail="contre-analyse momentanement indisponible")

    # 1er message : contexte Barzel + document + consignes initiales (turns[0]).
    # Tours suivants : l'historique metier tel quel (le doc reste dans le 1er message).
    mix_note = ("\n(Dossier potentiellement mixte : les donnees ci-dessus couvrent les cinq classes "
                "d'actifs de la ville.)" if asset_class == _CLS_ALL else "")
    msgs: list[dict] = [{
        "role": "user",
        "content": (f"{_context_for(asset_class, city)}{mix_note}\n\n"
                    f"# DOCUMENT EXTERNE\n{doc_text}\n\n"
                    f"# CONSIGNES DE L'INVESTISSEUR\n{turns[0].text}"),
    }]
    for turn in turns[1:]:
        msgs.append({"role": "assistant" if turn.role == "assistant" else "user", "content": turn.text})

    try:
        import anthropic

        client = anthropic.Anthropic(api_key=key)
        message = client.messages.create(
            model=_MODEL,
            max_tokens=_MAX_TOKENS,
            temperature=0.2,
            system=_system_for(city, lang),
            messages=msgs,
        )
        answer = "".join(b.text for b in message.content if getattr(b, "type", "") == "text").strip()
        return {"answer": strip_em_dashes(answer)}
    except HTTPException:
        raise
    except Exception as exc:  # noqa: BLE001 ; echec sobre, aucune fuite d'interne
        log.warning("second-opinion call failed: %s", type(exc).__name__)
        raise HTTPException(status_code=503, detail="contre-analyse momentanement indisponible")
